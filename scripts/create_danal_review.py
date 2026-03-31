#!/usr/bin/env python3
"""다날 검수요청서 자동 생성 - 스크린샷 캡처 + pptx 작성"""

import os
import time
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT_DIR = '/Users/taemin/Downloads/youtube-analysis'
SCREENSHOT_DIR = os.path.join(OUTPUT_DIR, '_screenshots')
os.makedirs(SCREENSHOT_DIR, exist_ok=True)

SITE_URL = 'https://www.bibllab.com'

def capture_screenshots():
    """사이트 스크린샷 캡처"""
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        context = browser.new_context(
            viewport={'width': 1280, 'height': 900},
            device_scale_factor=2
        )
        page = context.new_page()

        # 1. 메인 페이지
        print('1. 메인 페이지 캡처...')
        page.goto(SITE_URL, wait_until='networkidle')
        time.sleep(2)
        page.screenshot(path=os.path.join(SCREENSHOT_DIR, '01_main.png'), full_page=False)

        # 2. Footer 영역
        print('2. Footer 캡처...')
        page.evaluate('window.scrollTo(0, document.body.scrollHeight)')
        time.sleep(1)
        page.screenshot(path=os.path.join(SCREENSHOT_DIR, '02_footer.png'), full_page=False)

        # 3. 요금제 페이지
        print('3. 요금제 페이지 캡처...')
        page.goto(f'{SITE_URL}/pricing', wait_until='networkidle')
        time.sleep(2)
        page.screenshot(path=os.path.join(SCREENSHOT_DIR, '03_pricing.png'), full_page=False)

        # 4. 결제 페이지 (로그인 불필요한 화면까지만)
        print('4. 결제 페이지 캡처...')
        page.goto(f'{SITE_URL}/payment?plan=starter', wait_until='networkidle')
        time.sleep(2)
        page.screenshot(path=os.path.join(SCREENSHOT_DIR, '04_payment.png'), full_page=False)

        # 5. 서비스 소개 페이지
        print('5. 서비스 소개 페이지 캡처...')
        page.goto(f'{SITE_URL}/about', wait_until='networkidle')
        time.sleep(2)
        page.screenshot(path=os.path.join(SCREENSHOT_DIR, '05_about.png'), full_page=False)

        browser.close()
    print('스크린샷 캡처 완료!')


def build_pptx():
    """검수요청서 pptx 생성"""
    template_path = '/Users/taemin/Downloads/신용카드(일반) 검수요청서_양식.pptx'
    output_path = os.path.join(OUTPUT_DIR, '다날_검수요청서_세모골프.pptx')

    prs = Presentation(template_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # ── Slide 1: 표지 ──
    slide = prs.slides[0]
    # 기존 텍스트 상자 확인 및 표지 정보 추가
    # 표지에 텍스트 박스 추가
    left = Inches(1.5)
    top = Inches(2.0)
    width = Inches(7)
    height = Inches(0.5)

    items = [
        ('상 호 명', '세모골프 (bibl lab)'),
        ('서비스 URL', 'https://www.bibllab.com'),
        ('결제 프로세스', '요금제 선택 → 결제 페이지 → 신용카드 결제'),
        ('작 성 일 자', '2026. 03. 31'),
        ('작 성 자', '김태민 (대표)'),
    ]

    for i, (label, value) in enumerate(items):
        txBox = slide.shapes.add_textbox(left, top + Inches(i * 0.6), width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run_label = p.add_run()
        run_label.text = f'{label}: '
        run_label.font.size = Pt(14)
        run_label.font.bold = True
        run_label.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        run_value = p.add_run()
        run_value.text = value
        run_value.font.size = Pt(14)
        run_value.font.color.rgb = RGBColor(0x00, 0x00, 0x00)

    # ── Slide 2: 사이트 메인 ──
    slide = prs.slides[1]
    img_path = os.path.join(SCREENSHOT_DIR, '01_main.png')
    if os.path.exists(img_path):
        # URL 텍스트 추가
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f'URL: {SITE_URL}'
        run.font.size = Pt(10)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

        slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.0), Inches(9.4))

    # ── Slide 3: Footer 정보 ──
    slide = prs.slides[2]
    img_path = os.path.join(SCREENSHOT_DIR, '02_footer.png')
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0.3), Inches(1.0), Inches(9.4))

        # 필수 정보 체크리스트
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9), Inches(1.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = '✓ 상호명: 세모골프  ✓ 대표자명: 김태민  ✓ 사업자번호: 315-47-01018\n✓ 유선전화: 070-8027-2532  ✓ 주소: 경기도 수원시 권선구 세화로 151번길 29-2 1층\n✓ 통신판매업: 2023-수원권선-1549'
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── Slide 4: 결제 프로세스 ──
    slide = prs.slides[3]
    # TEST ID/PW 입력
    for shape in slide.shapes:
        if shape.has_text_frame:
            for p in shape.text_frame.paragraphs:
                if 'TEST ID' in p.text:
                    p.clear()
                    run = p.add_run()
                    run.text = 'TEST ID : 테스트 계정 없이 이용 가능 (Google 소셜 로그인)'
                    run.font.size = Pt(10)
                elif 'TEST PW' in p.text:
                    p.clear()
                    run = p.add_run()
                    run.text = 'TEST PW : Google 소셜 로그인 사용 (별도 비밀번호 없음)'
                    run.font.size = Pt(10)

    # 결제 프로세스 스크린샷 (요금제 + 결제 페이지)
    pricing_img = os.path.join(SCREENSHOT_DIR, '03_pricing.png')
    payment_img = os.path.join(SCREENSHOT_DIR, '04_payment.png')
    if os.path.exists(pricing_img):
        slide.shapes.add_picture(pricing_img, Inches(0.3), Inches(1.5), Inches(4.5))
    if os.path.exists(payment_img):
        slide.shapes.add_picture(payment_img, Inches(5.0), Inches(1.5), Inches(4.5))

    # 설명 텍스트
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '유튜브 데이터 분석 SaaS 서비스(bibl lab) 월 정기구독 결제입니다.\n요금제: Starter ₩49,000 / Pro ₩199,000 / Business ₩490,000 (월 구독형)'
    run.font.size = Pt(9)
    run.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    # ── Slide 5: 신용카드 결제창 ──
    # 이 슬라이드는 실제 카드 결제 팝업 스크린샷이 필요합니다.
    # 포트원 테스트 모듈 연동 후 캡처 필요
    slide = prs.slides[4]
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = '※ 포트원(PortOne) 결제 연동을 통한 신용카드 결제창입니다.\n다날 테스트 모듈 연동 완료 후 결제창 스크린샷을 업데이트 예정입니다.\n\n현재 KG이니시스/NHN KCP 테스트 모듈로 결제 연동이 완료되어 있으며,\n다날 채널키 발급 후 동일한 방식으로 연동됩니다.'
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # ── Slide 6-7: 카드사 인증창 ──
    # 다날 테스트 모듈 미연동 상태이므로 안내 텍스트
    for idx in [5, 6]:
        if idx < len(prs.slides):
            slide = prs.slides[idx]
            txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = '※ 다날 테스트 모듈 연동 후 카드사 인증창 스크린샷을 추가 제출 예정입니다.'
            run.font.size = Pt(12)
            run.font.color.rgb = RGBColor(0x66, 0x66, 0x66)

    # 환금성 슬라이드(9-12) 삭제 - 환금성 업종 아님
    # pptx에서 슬라이드 삭제는 복잡하므로 그대로 둠 (빈 상태 유지)

    prs.save(output_path)
    print(f'검수요청서 생성 완료: {output_path}')


if __name__ == '__main__':
    capture_screenshots()
    build_pptx()
